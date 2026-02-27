import json
import re
from typing import Any, Union

from langchain_core.messages import AIMessage
from langchain_core.runnables import RunnableLambda

from cuga.backend.activity_tracker.tracker import ActivityTracker
from cuga.backend.cuga_graph.nodes.shared.base_agent import BaseAgent
from cuga.backend.cuga_graph.state.agent_state import AgentState
from cuga.backend.cuga_graph.nodes.innovation.prompts.load_prompt import (
    InnovationEvaluationOutput,
    InnovationSummaryOutput,
    parser_evaluation,
    parser_summary,
)
from cuga.backend.llm.models import LLMManager
from cuga.backend.llm.utils.helpers import load_prompt_simple
from cuga.config import settings
from loguru import logger
import os
import tempfile
import shutil
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from gpt_researcher import GPTResearcher
except ImportError:
    GPTResearcher = None
    logger.warning("gpt-researcher not installed. ResearchAgent will fail.")

tracker = ActivityTracker()
llm_manager = LLMManager()

# URL pattern for detecting URLs in user input
URL_PATTERN = re.compile(r'https?://[^\s]+')


class ResearchAgent(BaseAgent):
    """General-purpose research agent powered by GPT Researcher.
    
    Supports 4 modes:
    - research: Web-only research for queries without documents/URLs
    - url_summary: Summarize content from a URL
    - summary: Summarize uploaded documents
    - evaluation: Innovation evaluation with Search-X scoring
    """

    def __init__(self):
        super().__init__()
        self.name = "ResearchAgent"
        
        # Use a capable model for evaluation logic
        dyna_model = settings.agent.plan_controller.model
        llm = llm_manager.get_model(dyna_model)
        
        # --- Evaluation Chain ---
        prompt_eval = load_prompt_simple(
            "./prompts/system.jinja2",
            "./prompts/user_evaluation.jinja2",
            model_config=dyna_model,
            format_instructions=BaseAgent.get_format_instructions(parser_evaluation),
        )
        output_parser_eval = RunnableLambda(lambda x: ResearchAgent.output_parser(x, self.name))
        
        self.chain_evaluation = BaseAgent.get_chain(prompt_eval, llm, InnovationEvaluationOutput) | (
            output_parser_eval
        )

        # --- Summarization Chain ---
        prompt_summary = load_prompt_simple(
            "./prompts/system_summary.jinja2",
            "./prompts/user_summary.jinja2",
            model_config=dyna_model,
            format_instructions=BaseAgent.get_format_instructions(parser_summary),
        )
        output_parser_summary = RunnableLambda(lambda x: ResearchAgent.output_parser(x, self.name))
        
        self.chain_summary = BaseAgent.get_chain(prompt_summary, llm, InnovationSummaryOutput) | (
            output_parser_summary
        )

    def _extract_text_from_file(self, file_path: str) -> str:
        """Extract text from PDF, DOCX, or PPTX files for safer processing."""
        ext = os.path.splitext(file_path)[1].lower()
        text = ""
        
        try:
            if ext == ".pptx" and Presentation:
                logger.info(f"Attempting to extract text from PPTX: {file_path}")
                try:
                    prs = Presentation(file_path)
                    for slide_num, slide in enumerate(prs.slides, 1):
                        try:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    text += shape.text + "\n"
                        except Exception as shape_error:
                            logger.warning(f"Error processing slide {slide_num}: {shape_error}")
                            continue
                    
                    if text.strip():
                        logger.info(f"Successfully extracted {len(text)} characters from PPTX")
                        return text
                    else:
                        logger.warning("No text extracted from PPTX, will let GPT Researcher handle it")
                        return ""
                        
                except Exception as pptx_error:
                    logger.warning(f"PPTX extraction failed with error: {pptx_error}. Will let GPT Researcher handle the file directly.")
                    return ""
            
            elif ext == ".docx" and Document:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                return text
            
            elif ext == ".txt":
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            # For others (PDF etc.), we let GPTResearcher handle them or
            # we could add more loaders here if needed.
            return ""
        except Exception as e:
            logger.warning(f"Failed to extract text from {file_path}: {e}. Will let GPT Researcher handle it.")
            return ""

    @staticmethod
    def detect_mode(user_input: str) -> str:
        """Determine the research mode from user input.
        
        Returns one of: 'research', 'url_summary', 'summary', 'evaluation'
        """
        lower_input = user_input.lower()
        has_url = bool(URL_PATTERN.search(user_input))
        
        # Check for evaluation intent first (most specific)
        eval_keywords = [
            "evaluate", "evaluation", "disclosure", "novelty", "patent", "prior art",
            "score", "classify", "search-1", "search-2", "search-3", "publish", "assess", "rating"
        ]
        if any(kw in lower_input for kw in eval_keywords):
            return "evaluation"
        
        # URL summarization
        summary_keywords = ["summarize", "summary", "overview", "explain", "what is this"]
        if has_url and any(kw in lower_input for kw in summary_keywords):
            return "url_summary"
        
        # URL without summary keyword — treat as research about that URL
        if has_url:
            return "url_summary"
        
        # Document summarization (file path present)
        if any(kw in lower_input for kw in summary_keywords):
            return "summary"
        
        # Default: web research
        return "research"

    @staticmethod
    def output_parser(result: Union[InnovationEvaluationOutput, InnovationSummaryOutput], name) -> AIMessage:
        logger.debug(
            f"\n\n------\n\nResearch Output: {json.dumps(result.model_dump(), indent=4)} \n\n------\n\n"
        )
        return AIMessage(content=json.dumps(result.model_dump()), name=name)

    async def run(self, state: AgentState) -> AIMessage:
        if not GPTResearcher:
            return AIMessage(content=json.dumps({"error": "gpt-researcher not installed"}), name=self.name)

        user_input = state.input
        mode = self.detect_mode(user_input)
        logger.info(f"ResearchAgent: Detected mode='{mode}' for input: {user_input[:100]}")

        doc_path = os.environ.get("DOC_PATH")
        # ------ Build query and report_source based on mode ------
        if mode == "research":
            # Web-only research, no documents
            query = user_input
            report_source = "web"
            
        elif mode == "url_summary":
            # Extract URL from input
            urls = URL_PATTERN.findall(user_input)
            url = urls[0] if urls else user_input
            query = f"Summarize and analyze the content at: {url}"
            report_source = "web"
            
        elif mode == "summary":
            doc_path = os.environ.get("DOC_PATH")
            if not doc_path:
                logger.warning("DOC_PATH env var not set for summary mode.")
            query = f"Summarize the key features, claims, and applications of the content in the uploaded documents. User intent: {user_input}"
            report_source = "hybrid"
            
        elif mode == "evaluation":
            doc_path = os.environ.get("DOC_PATH")
            if not doc_path:
                logger.warning("DOC_PATH env var not set for evaluation mode.")
            query = f"Analyze the novelty, key features, and prior art for the innovation described in the uploaded documents. User intent: {user_input}"
            report_source = "hybrid"
        else:
            query = user_input
            report_source = "web"

        # ------ Handle DOC_PATH for hybrid/local modes ------
        actual_doc_path = doc_path
        temp_dir = None
        if report_source == "hybrid" and doc_path:
            # Create a temp directory for GPT Researcher to walk
            temp_dir = tempfile.mkdtemp()
            actual_doc_path = temp_dir
            
            if os.path.isfile(doc_path):
                # Single file: For PPTX files with memory issues, just copy them directly
                # GPT Researcher has its own robust document loaders
                ext = os.path.splitext(doc_path)[1].lower()
                
                # Pre-extract text for DOCX and PPTX to avoid GPT Researcher loader issues
                if ext in [".docx", ".pptx"]:
                    logger.info(f"ResearchAgent: Pre-extracting text from {ext.upper()} {doc_path}")
                    extracted_text = self._extract_text_from_file(doc_path)
                    if extracted_text:
                        txt_path = os.path.join(temp_dir, os.path.basename(doc_path) + ".txt")
                        with open(txt_path, "w", encoding="utf-8") as f:
                            f.write(extracted_text)
                    else:
                        # Fallback to copy if extraction failed
                        shutil.copy(doc_path, os.path.join(temp_dir, os.path.basename(doc_path)))
                else:
                    shutil.copy(doc_path, os.path.join(temp_dir, os.path.basename(doc_path)))

                logger.info(f"ResearchAgent: Prepared document in temp directory {temp_dir}")
            elif os.path.isdir(doc_path):
                # If it's a directory, we'll process all files inside it
                for item in os.listdir(doc_path):
                    item_path = os.path.join(doc_path, item)
                    if os.path.isfile(item_path):
                        ext = os.path.splitext(item_path)[1].lower()
                        if ext in [".pptx", ".docx"]:
                            extracted_text = self._extract_text_from_file(item_path)
                            if extracted_text:
                                txt_path = os.path.join(temp_dir, item + ".txt")
                                with open(txt_path, "w", encoding="utf-8") as f:
                                    f.write(extracted_text)
                                continue
                        shutil.copy(item_path, os.path.join(temp_dir, item))
            elif not os.path.exists(doc_path):
                logger.warning(f"ResearchAgent: DOC_PATH {doc_path} does not exist.")

        try:
            # Set environment variable temporarily for gpt-researcher if needed
            # (Some versions rely on DOC_PATH env var)
            if actual_doc_path:
                os.environ["DOC_PATH"] = actual_doc_path
                
            # Ensure team-allowed models are used to avoid defaulting to o4-mini
            if not os.environ.get("STRATEGIC_LLM"):
                os.environ["STRATEGIC_LLM"] = "Azure/gpt-4.1"
                
            logger.info(f"Running GPT Researcher: mode={mode}, report_source={report_source}, docs={actual_doc_path}")
            logger.info(f"Research Query: {query}")
            
            # Initialize researcher with appropriate report type
            report_type = "research_report"
            # If user asks for deep research explicitly or query is complex, we could use "deep"
            # For now, stick to standard "research_report" which is better than default "custom_report"
            
            researcher = GPTResearcher(
                query=query, 
                report_source=report_source, 
                report_type=report_type,
                documents=actual_doc_path if report_source == "hybrid" else None
            )
            
            logger.info("Conducting research...")
            await researcher.conduct_research()
            
            logger.info("Writing report...")
            report = await researcher.write_report()
            
            logger.info(f"GPT Researcher completed successfully. Report length: {len(report)}")
            logger.debug(f"Research Report Content: {report[:500]}...")

            # ------ Process result based on mode ------
            if mode in ("research", "url_summary"):
                # Return raw markdown report — no structured output parser needed
                return AIMessage(content=json.dumps({"report": report, "mode": mode}), name=self.name)
            
            elif mode == "summary":
                chain_input = {"research_report": report}
                result = await self.chain_summary.ainvoke(chain_input)
                return result
            
            elif mode == "evaluation":
                # --- STEP 1: Understanding ---
                # Use the summary chain to extract structured understanding from the report
                # (The report contains the researcher's analysis of the local documents)
                logger.info("Evaluation Mode: Step 1 - Extracting Understanding")
                summary_input = {"research_report": report}
                summary_msg = await self.chain_summary.ainvoke(summary_input)
                innovation_understanding = summary_msg.content # This is a JSON string of InnovationSummaryOutput
                
                # --- STEP 2: Judging ---
                # Feed both Understanding and Research Report into the evaluation chain
                logger.info("Evaluation Mode: Step 2 - Professional Judging")
                eval_input = {
                    "innovation_understanding": innovation_understanding,
                    "research_report": report
                }
                logger.info(f"Invoking evaluation chain with keys: {list(eval_input.keys())}")
                result = await self.chain_evaluation.ainvoke(eval_input)
                
                # Append the full research report for transparency as requested
                try:
                    content_dict = json.loads(result.content)
                    content_dict["research_report"] = report
                    result.content = json.dumps(content_dict)
                except Exception as e:
                    logger.warning(f"Failed to append research_report to JSON output: {e}")
                
                return result

        except Exception as e:
            logger.error(f"Error in ResearchAgent: {e}")
            return AIMessage(content=json.dumps({"error": str(e)}), name=self.name)
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
                logger.debug(f"ResearchAgent: Cleaned up temp directory {temp_dir}")
                
            # Clean up the DOC_PATH environment variable if we set it temporarily
            if actual_doc_path and "DOC_PATH" in os.environ and os.environ["DOC_PATH"] == actual_doc_path:
                del os.environ["DOC_PATH"]
                logger.debug(f"ResearchAgent: Cleaned up DOC_PATH environment variable")
